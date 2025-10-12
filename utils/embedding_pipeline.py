import os
import fitz
import openai
from docx import Document as DocxDocument
from pptx import Presentation
from supabase import create_client

CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
VECTOR_DIM = 3072
EMBED_MODEL = "text-embedding-3-large"

def create_supabase_client():
    from dotenv import load_dotenv
    load_dotenv()
    url = os.getenv("SUPABASE_URL")
    key = os.getenv("SUPABASE_SERVICE_KEY")
    return create_client(url, key)

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    else:
        return ""

def parse_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text("text") for page in doc)

def parse_docx(path):
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)

def parse_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)

def chunk_text(text, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    words = text.split()
    step = size - overlap
    for i in range(0, len(words), step):
        yield " ".join(words[i:i+size])
        if i + size >= len(words):
            break

def embed_text(text):
    resp = openai.embeddings.create(model=EMBED_MODEL, input=[text])
    emb = resp.data[0].embedding
    if len(emb) != VECTOR_DIM:
        raise ValueError(f"Expected {VECTOR_DIM} dims, got {len(emb)}")
    return emb

def embed_and_upload(file_path):
    supabase = create_supabase_client()
    text = extract_text(file_path)
    if not text.strip():
        return {"error": "No text extracted"}

    chunks = list(chunk_text(text))
    for i, chunk in enumerate(chunks):
        emb = embed_text(chunk)
        meta = {
            "source_file": os.path.basename(file_path),
            "chunk_index": i,
            "file_type": os.path.splitext(file_path)[1].lower(),
        }
        supabase.table("documents").insert({
            "content": chunk,
            "embedding": emb,
            "metadata": meta
        }).execute()
    return {"message": f"{len(chunks)} chunks embedded and uploaded"}
