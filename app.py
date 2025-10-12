import os, sys, tempfile, traceback, re
from zipfile import ZipFile
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from supabase import create_client
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from openai import OpenAI

# -------- logging --------
os.environ["PYTHONUNBUFFERED"] = "1"
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(line_buffering=True)
def log(msg: str): print(msg, flush=True)

# -------- env --------
load_dotenv()
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

EMBED_MODEL = "text-embedding-3-large"
CHAT_MODEL  = "gpt-4o-mini"
VECTOR_DIM  = 3072
TOP_K = 10
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
BUCKET = "materials"

# -------- clients --------
app = Flask(__name__, static_folder="static", static_url_path="/static")
CORS(app)
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
oai = OpenAI(api_key=OPENAI_API_KEY)

# -------- system prompt --------
def load_system_prompt(filepath="system_message.txt"):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception:
        return ("You are a helpful assistant that answers questions strictly using the provided context. "
                "If the context doesn't include enough information, say so clearly.")
SYSTEM_PROMPT = load_system_prompt()

# ==============================
# Text extraction
# ==============================
def extract_pdf_text(path: str) -> str:
    doc = fitz.open(path)
    parts = [page.get_text("text") for page in doc if page.get_text("text")]
    return "\n".join(parts)

def extract_docx_text(path: str) -> str:
    out = []
    try:
        doc = DocxDocument(path)
        for para in doc.paragraphs:
            t = para.text.strip()
            if t: out.append(t)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells: out.append(" | ".join(cells))
    except Exception as e:
        log(f"❌ python-docx failed: {e}")
        traceback.print_exc()
    # Fallback XML
    try:
        with ZipFile(path, "r") as z:
            xml_files = [f for f in z.namelist() if f.startswith("word/") and f.endswith(".xml")]
            combined = ""
            for f in xml_files:
                data = z.read(f).decode("utf-8", errors="ignore")
                combined += " " + re.sub(r"<[^>]+>", " ", data)
            combined = re.sub(r"\s+", " ", combined)
            if len(combined.strip()) > len("\n".join(out)):
                log(f"⚙️ deep XML fallback extracted {len(combined)} chars")
                return combined
    except Exception as e:
        log(f"❌ deep fallback failed: {e}")
    return "\n".join(out)

def extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line: out.append(line)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells: out.append(" | ".join(cells))
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t: out.append(t)
    return "\n".join(out)

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":  return extract_pdf_text(path)
        if ext == ".docx": return extract_docx_text(path)
        if ext == ".pptx": return extract_pptx_text(path)
        return ""
    except Exception as e:
        log(f"❌ extract_text failed for {path}: {e}")
        return ""

# ==============================
# Chunking + Embedding
# ==============================
def chunk_text(text: str, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    words = text.split()
    if not words: return
    if len(words) <= chunk_size:
        yield " ".join(words); return
    step = max(1, chunk_size - overlap)
    for i in range(0, len(words), step):
        chunk = words[i:i+chunk_size]
        yield " ".join(chunk)
        if i + chunk_size >= len(words): break

def embed_text(text: str):
    resp = oai.embeddings.create(model=EMBED_MODEL, input=[text])
    emb = resp.data[0].embedding
    if len(emb) != VECTOR_DIM:
        log(f"⚠️ embedding dims {len(emb)} != expected {VECTOR_DIM}")
    return emb

def process_and_store(path: str, filename: str):
    log(f"⚙️ process_and_store for {filename}")
    text = extract_text(path)
    log(f"📄 Extracted {len(text)} characters")
    if not text.strip(): return {"status": "no_text"}

    chunks = list(chunk_text(text))
    total = 0
    for i, chunk in enumerate(chunks):
        try:
            emb = embed_text(chunk)
            supabase.table("documents").insert({
                "content": chunk,
                "embedding": emb,
                "metadata": {
                    "source_file": filename,
                    "chunk_index": i,
                    "file_type": os.path.splitext(filename)[1].lower()
                }
            }).execute()
            total += 1
        except Exception as e:
            log(f"❌ insert chunk {i}: {e}")
    log(f"✅ embedded {total} chunks for {filename}")
    return {"status": "success", "chunks": total}

# ==============================
# Routes
# ==============================
@app.route("/")
def index():
    return send_from_directory(".", "index.html")

@app.route("/manage")
def manage():
    return send_from_directory("templates", "manage.html")

@app.route("/api/files")
def list_files():
    """List files from the files table (authoritative)."""
    try:
        data = supabase.table("files").select("*").order("uploaded_at", desc=True).execute()
        return jsonify(data.data or [])
    except Exception as e:
        log(f"❌ list_files error: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/upload", methods=["POST"])
def upload_file():
    """Upload to Storage, insert files row, then embed."""
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No file provided"}), 400

    name = secure_filename(file.filename)
    data = file.read()
    log(f"✅ Upload received: {name} | {len(data)} bytes")

    tmp_path = os.path.join(tempfile.gettempdir(), name)
    with open(tmp_path, "wb") as f:
        f.write(data)

    # 1) Storage upload (flat, upsert)
    try:
        res = supabase.storage.from_(BUCKET).upload(name, data, upsert=True)
        if hasattr(res, "error") and res.error is not None:
            raise Exception(res.error.message)
        log(f"📦 Stored: {name}")
    except Exception as e:
        log(f"❌ upload error: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

    # 2) Insert into files table (links UI <-> Storage)
    try:
        supabase.table("files").insert({
            "file_name": name,
            "file_type": os.path.splitext(name)[1].lower(),
            "source_path": name
            # uploaded_at handled by DB default NOW() if present
        }).execute()
        log(f"📝 files row inserted: {name}")
    except Exception as e:
        # Non-fatal: UI should still see storage, but we keep behavior consistent
        log(f"⚠️ files insert failed: {e}")

    # 3) Embed into documents (non-fatal)
    try:
        result = process_and_store(tmp_path, name)
        log(f"🧠 Embedding result: {result}")
    except Exception as e:
        log(f"⚠️ embedding failed: {e}")
        result = {"status": "embedding_failed", "error": str(e)}

    return jsonify({
        "message": f"{name} uploaded",
        "file_name": name,
        "file_type": os.path.splitext(name)[1],
        "embedding_result": result
    })

@app.route("/delete/<name>", methods=["DELETE"])
def delete_file(name):
    """Delete from Storage, files, and documents."""
    try:
        # Storage
        res = supabase.storage.from_(BUCKET).remove([name])
        if hasattr(res, "error") and res.error is not None:
            raise Exception(res.error.message)
        # files table
        supabase.table("files").delete().eq("file_name", name).execute()
        # embeddings
        supabase.table("documents").delete().filter("metadata->>source_file", "eq", name).execute()
        log(f"🗑️ Deleted {name} from storage, files, documents")
        return jsonify({"message": f"{name} deleted"})
    except Exception as e:
        log(f"❌ delete error: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/chat", methods=["POST"])
def chat():
    user_input = request.json.get("message", "").strip()
    if not user_input:
        return jsonify({"error": "No input provided"}), 400
    try:
        emb = oai.embeddings.create(model=EMBED_MODEL, input=[user_input]).data[0].embedding
        resp = supabase.rpc("match_documents", {"query_embedding": emb, "match_count": TOP_K}).execute()
        matches = getattr(resp, "data", []) or []
        if not matches:
            return jsonify({"response": "No relevant materials.\n**Sources:** None"})
        blocks, srcs = [], set()
        for m in matches:
            c = m.get("content", "").strip()
            s = m.get("source_file", "Unknown")
            if c:
                blocks.append(f"[source_file: {s}]\n{c}")
                srcs.add(s)
        context = "\n\n".join(blocks)
        src_txt = "**Sources:** " + ", ".join(f"`{s}`" for s in sorted(srcs))
        ans = oai.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": f"Answer using only this context:\n\n{context}\n\nQuestion: {user_input}"}
            ],
            temperature=0.2
        ).choices[0].message.content.strip()
        return jsonify({"response": f"{ans}\n\n{src_txt}"})
    except Exception as e:
        log(f"❌ chat error: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# -------- run --------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
