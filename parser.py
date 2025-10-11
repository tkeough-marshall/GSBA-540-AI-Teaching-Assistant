# app/parser.py
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import openai
import io

CHUNK_SIZE = 500
CHUNK_OVERLAP = 100

def extract_text(file_bytes: bytes, ext: str) -> str:
    if ext == ".pdf":
        return parse_pdf(file_bytes)
    elif ext == ".docx":
        return parse_docx(file_bytes)
    elif ext == ".pptx":
        return parse_pptx(file_bytes)
    return ""

def parse_pdf(file_bytes: bytes) -> str:
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        return "\n".join([page.get_text("text") for page in doc])

def parse_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text)

def parse_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    texts = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text]
    return "\n".join(texts)

def chunk_text(text: str, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    words = text.split()
    step = chunk_size - overlap
    for i in range(0, len(words), step):
        yield " ".join(words[i:i + chunk_size])
        if i + chunk_size >= len(words):
            break

def embed_text(text: str):
    resp = openai.embeddings.create(input=[text], model="text-embedding-3-large")
    return resp.data[0].embedding