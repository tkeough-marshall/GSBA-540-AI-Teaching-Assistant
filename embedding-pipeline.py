# ===============================
# 📌 CONFIGURATION (EDIT THESE)
# ===============================
CHUNK_SIZE = 500           # number of words per chunk (adjust if needed)
CHUNK_OVERLAP = 100        # overlap between successive chunks
VECTOR_DIM = 3072          # embedding dimension you're using
EMBED_MODEL = "text-embedding-3-large"  # adjust if using a different model
ALLOWED_EXTENSIONS = (".pdf", ".docx", ".pptx")
MAX_FILES = None             # number of files to process for testing (None = process all)

# ===============================
# 🔧 IMPORTS & SETUP
# ===============================
import os
import fitz  # PyMuPDF for PDF parsing
import openai
from docx import Document as DocxDocument
from pptx import Presentation
from tqdm import tqdm
from tkinter import Tk, filedialog
from supabase import create_client
from dotenv import load_dotenv

load_dotenv()

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if not SUPABASE_URL or not SUPABASE_KEY or not OPENAI_API_KEY:
    raise RuntimeError("Missing one or more environment variables: SUPABASE_URL, SUPABASE_SERVICE_KEY, OPENAI_API_KEY")

openai.api_key = OPENAI_API_KEY
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

# ===============================
# 📄 FILE PARSERS
# ===============================
def parse_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text_pages = []
    for page in doc:
        txt = page.get_text("text")
        if txt:
            text_pages.append(txt)
    return "\n".join(text_pages)

def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text]
    return "\n".join(paragraphs)

def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)
    return "\n".join(texts)

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    else:
        return ""

# ===============================
# 🔢 CHUNKING + EMBEDDING
# ===============================
def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
    words = text.split()
    if len(words) <= chunk_size:
        yield " ".join(words)
        return
    step = chunk_size - overlap
    for i in range(0, len(words), step):
        chunk = words[i : i + chunk_size]
        yield " ".join(chunk)
        if i + chunk_size >= len(words):
            break

def embed_text(text: str):
    resp = openai.embeddings.create(
        input=[text],
        model=EMBED_MODEL
    )
    embedding = resp.data[0].embedding
    if len(embedding) != VECTOR_DIM:
        raise ValueError(f"Embedding returned dimension {len(embedding)} but expected {VECTOR_DIM}")
    return embedding

# ===============================
# 🚀 UPLOAD TO SUPABASE
# ===============================
def upload_to_supabase(content: str, file_path: str, chunk_idx: int):
    embedding = embed_text(content)
    metadata = {
        "source_file": os.path.basename(file_path),
        "chunk_index": chunk_idx,
        "file_type": os.path.splitext(file_path)[1].lower(),
    }
    result = supabase.table("documents").insert({
        "content": content,
        "embedding": embedding,
        "metadata": metadata
    }).execute()
    if result.get("error"):
        print("❌ Supabase insert error:", result["error"])

# ===============================
# 🧠 FILE PROCESSING
# ===============================
def process_file(file_path: str):
    print(f"📘 Processing: {file_path}")
    raw = extract_text(file_path)
    if not raw or not raw.strip():
        print(f"⚠️ Skipping (no text): {file_path}")
        return

    chunks = list(chunk_text(raw))
    with tqdm(total=len(chunks), desc=f"Chunks for {os.path.basename(file_path)}", leave=False) as pbar:
        for idx, chunk in enumerate(chunks):
            try:
                upload_to_supabase(chunk, file_path, idx)
            except Exception as e:
                print(f"❌ Error chunk {idx} of {file_path}: {e}")
            pbar.update(1)

def select_folder() -> str:
    root = Tk()
    root.withdraw()
    return filedialog.askdirectory(title="Select Course Files Folder")

# ===============================
# 🎯 MAIN
# ===============================
def main():
    folder = select_folder()
    if not folder:
        print("❌ No folder selected, exiting.")
        return

    all_files = [
        os.path.join(folder, fname)
        for fname in os.listdir(folder)
        if fname.lower().endswith(ALLOWED_EXTENSIONS)
    ]
    files = all_files[:MAX_FILES] if MAX_FILES is not None else all_files

    print(f"\n🚀 Starting ingestion of {len(files)} files from folder: {folder}\n")

    for fp in tqdm(files, desc="Processing files"):
        try:
            process_file(fp)
        except Exception as e:
            print(f"❌ Failed file {fp}: {e}")

    print("\n✅ All done! Ingestion complete.\n")

if __name__ == "__main__":
    main()
